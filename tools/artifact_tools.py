"""Cross-platform local tools for common document artifact plugins."""

from __future__ import annotations

import html
import re
from pathlib import Path
from typing import Any

from tools.file_tools import FileTools


class ArtifactTools:
    """Read and create office artifacts while staying inside the project."""

    def __init__(self, project_root: str | Path):
        self.files = FileTools(project_root)

    def _path(self, path: str, suffix: str) -> Path:
        resolved = self.files._resolve_path(path)
        if resolved.suffix.lower() != suffix:
            raise ValueError(f"Expected a {suffix} file: {path}")
        return resolved

    def _result(self, path: Path, status: str = "written", **extra) -> dict:
        return {
            "status": status,
            "path": self.files._relative_path(path),
            "absolute_path": str(path),
            "size": path.stat().st_size if path.exists() else 0,
            **extra,
        }

    async def read_docx(self, path: str, start_block: int = 1,
                        max_blocks: int = 80, max_chars: int = 16_000,
                        **kwargs) -> dict:
        try:
            from docx import Document
        except ImportError:
            return self._dependency_error("python-docx")
        resolved = self._path(path, ".docx")
        if not resolved.is_file():
            return {"status": "error", "error": f"File not found: {path}"}
        document = Document(str(resolved))
        blocks = [paragraph.text for paragraph in document.paragraphs]
        for table_index, table in enumerate(document.tables, 1):
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                blocks.append(
                    f"[Table {table_index}] " + " | ".join(cells)
                )
        start = max(1, int(start_block or 1))
        count = max(1, min(200, int(max_blocks or 80)))
        limit = max(2_000, min(32_000, int(max_chars or 16_000)))
        selected = blocks[start - 1:start - 1 + count]
        content = "\n\n".join(selected)[:limit]
        end = min(len(blocks), start - 1 + len(selected))
        return self._result(
            resolved, "success", content=content, total_blocks=len(blocks),
            start_block=start, end_block=end, has_more=end < len(blocks),
            next_block=end + 1 if end < len(blocks) else None,
        )

    async def write_docx(self, path: str, content: str, title: str = "",
                         **kwargs) -> dict:
        try:
            from docx import Document
        except ImportError:
            return self._dependency_error("python-docx")
        resolved = self._path(path, ".docx")
        resolved.parent.mkdir(parents=True, exist_ok=True)
        document = Document()
        if title.strip():
            document.add_heading(title.strip(), level=0)
        for kind, text, level in self._markdown_blocks(content):
            if kind == "heading":
                document.add_heading(text, level=min(9, max(1, level)))
            elif kind == "bullet":
                document.add_paragraph(text, style="List Bullet")
            elif kind == "number":
                document.add_paragraph(text, style="List Number")
            else:
                document.add_paragraph(text)
        document.save(str(resolved))
        return self._result(resolved)

    async def read_pptx(self, path: str, start_slide: int = 1,
                        max_slides: int = 10, max_chars: int = 16_000,
                        **kwargs) -> dict:
        try:
            from pptx import Presentation
        except ImportError:
            return self._dependency_error("python-pptx")
        resolved = self._path(path, ".pptx")
        if not resolved.is_file():
            return {"status": "error", "error": f"File not found: {path}"}
        presentation = Presentation(str(resolved))
        total = len(presentation.slides)
        start = max(1, int(start_slide or 1))
        count = max(1, min(30, int(max_slides or 10)))
        chunks: list[str] = []
        end = start - 1
        for number in range(start, min(total, start + count - 1) + 1):
            slide = presentation.slides[number - 1]
            texts = [
                str(shape.text).strip() for shape in slide.shapes
                if getattr(shape, "has_text_frame", False)
                and str(getattr(shape, "text", "")).strip()
            ]
            chunks.append(f"--- Slide {number} ---\n" + "\n".join(texts))
            end = number
        limit = max(2_000, min(32_000, int(max_chars or 16_000)))
        return self._result(
            resolved, "success", content="\n\n".join(chunks)[:limit],
            slide_count=total, start_slide=start, end_slide=end,
            has_more=end < total, next_slide=end + 1 if end < total else None,
        )

    async def write_pptx(self, path: str, slides: list[dict],
                         title: str = "", subtitle: str = "", **kwargs) -> dict:
        try:
            from pptx import Presentation
        except ImportError:
            return self._dependency_error("python-pptx")
        if not isinstance(slides, list) or not slides:
            return {"status": "error", "error": "slides must be a non-empty array"}
        resolved = self._path(path, ".pptx")
        resolved.parent.mkdir(parents=True, exist_ok=True)
        presentation = Presentation()
        if title.strip():
            cover = presentation.slides.add_slide(presentation.slide_layouts[0])
            cover.shapes.title.text = title.strip()
            if len(cover.placeholders) > 1:
                cover.placeholders[1].text = subtitle.strip()
        for raw_slide in slides:
            if not isinstance(raw_slide, dict):
                continue
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = str(raw_slide.get("title") or "")
            body = slide.placeholders[1].text_frame
            body.clear()
            bullets = raw_slide.get("bullets") or raw_slide.get("content") or []
            if isinstance(bullets, str):
                bullets = [line for line in bullets.splitlines() if line.strip()]
            for index, bullet in enumerate(bullets):
                paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
                if isinstance(bullet, dict):
                    paragraph.text = str(bullet.get("text") or "")
                    paragraph.level = max(0, min(4, int(bullet.get("level") or 0)))
                else:
                    paragraph.text = str(bullet)
        presentation.save(str(resolved))
        return self._result(resolved, slide_count=len(presentation.slides))

    async def write_pdf(self, path: str, content: str, title: str = "",
                        **kwargs) -> dict:
        try:
            from reportlab.lib.enums import TA_LEFT
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
            from reportlab.lib.units import mm
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.cidfonts import UnicodeCIDFont
            from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
        except ImportError:
            return self._dependency_error("reportlab")
        resolved = self._path(path, ".pdf")
        resolved.parent.mkdir(parents=True, exist_ok=True)
        font_name = "STSong-Light"
        try:
            pdfmetrics.registerFont(UnicodeCIDFont(font_name))
        except Exception:
            font_name = "Helvetica"
        styles = getSampleStyleSheet()
        normal = ParagraphStyle(
            "RockCoreBody", parent=styles["BodyText"], fontName=font_name,
            fontSize=10.5, leading=17, alignment=TA_LEFT,
            spaceAfter=4 * mm, wordWrap="CJK",
        )
        heading = ParagraphStyle(
            "RockCoreHeading", parent=normal, fontSize=16, leading=22,
            spaceBefore=4 * mm, spaceAfter=3 * mm,
        )
        story: list[Any] = []
        if title.strip():
            story.extend([Paragraph(html.escape(title.strip()), heading), Spacer(1, 2 * mm)])
        for kind, text, level in self._markdown_blocks(content):
            prefix = "• " if kind == "bullet" else ""
            style = heading if kind == "heading" else normal
            story.append(Paragraph(html.escape(prefix + text), style))
        document = SimpleDocTemplate(
            str(resolved), pagesize=A4, leftMargin=18 * mm,
            rightMargin=18 * mm, topMargin=18 * mm, bottomMargin=18 * mm,
            title=title.strip(),
        )
        document.build(story)
        return self._result(resolved)

    @staticmethod
    def _markdown_blocks(content: str):
        for raw_line in str(content or "").splitlines():
            line = raw_line.strip()
            if not line:
                continue
            heading = re.match(r"^(#{1,6})\s+(.+)$", line)
            if heading:
                yield "heading", heading.group(2).strip(), len(heading.group(1))
            elif re.match(r"^[-*+]\s+", line):
                yield "bullet", re.sub(r"^[-*+]\s+", "", line), 0
            elif re.match(r"^\d+[.)]\s+", line):
                yield "number", re.sub(r"^\d+[.)]\s+", "", line), 0
            else:
                yield "paragraph", line, 0

    @staticmethod
    def _dependency_error(package: str) -> dict:
        return {
            "status": "dependency_missing",
            "error": f"Required packaged dependency is missing: {package}",
        }
